from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

OUT = Path('/usr/share/nginx/heratio/Trust-by-Design_Heratio_Show-and-Tell_20min_DRAFT.pptx')
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(11, 35, 49); TEAL = RGBColor(15, 118, 110); MINT = RGBColor(213, 245, 239)
GOLD = RGBColor(226, 169, 62); WHITE = RGBColor(255,255,255); INK = RGBColor(26,40,48)
PALE = RGBColor(244,248,247); GREY = RGBColor(94,108,115); RED = RGBColor(187,55,45)
ASSET = Path('/usr/share/nginx/conferences/DSAC - Archives and Records in Transition')
DSAC = ASSET / '_assets/DSAC.jpg'
CONTROL_DIAGRAM = ASSET / 'presentation/control_points_diagram.png'
CONFIDENCE_CHART = ASSET / 'presentation/fr_confidence_chart.png'
DOC_A = ASSET / 'pilot/images/0003.jpg'
DOC_B = ASSET / 'pilot/images/0009.jpg'
DOC_C = ASSET / 'pilot/images/0011.jpeg'

def box(slide, x,y,w,h, fill=WHITE, line=RGBColor(214,224,226), radius=True):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line
    return sh

def text(slide, value, x,y,w,h, size=24, color=INK, bold=False, align=PP_ALIGN.LEFT, font='Aptos', margin=.08):
    tb=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf=tb.text_frame
    tf.clear(); tf.margin_left=tf.margin_right=Inches(margin); tf.margin_top=tf.margin_bottom=Inches(margin)
    tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.text=value; p.alignment=align
    p.font.name=font; p.font.size=Pt(size); p.font.bold=bold; p.font.color.rgb=color
    return tb

def picture(slide, path, x, y, w, h, border=WHITE):
    with Image.open(path) as im:
        iw, ih = im.size
    target = w / h; source = iw / ih
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if source > target:
        crop = (1 - target/source) / 2
        pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - source/target) / 2
        pic.crop_top = crop; pic.crop_bottom = crop
    pic.line.color.rgb = border; pic.line.width = Pt(1.5)
    return pic

def base(title, kicker=None, dark=False):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg=s.background.fill; bg.solid(); bg.fore_color.rgb=NAVY if dark else PALE
    if kicker: text(s,kicker.upper(),.6,.28,8,.3,10,GOLD if dark else TEAL,True)
    text(s,title,.6,.66,12.1,.75,30,WHITE if dark else NAVY,True)
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(.6),Inches(1.53),Inches(1.15),Inches(.06)).fill.solid(); s.shapes[-1].fill.fore_color.rgb=GOLD if dark else TEAL; s.shapes[-1].line.fill.background()
    return s

def footer(s,n,time):
    text(s,f'{n:02d}',12.25,7.03,.45,.24,9,WHITE if s.background.fill.fore_color.rgb==NAVY else GREY,True,PP_ALIGN.RIGHT)
    text(s,time,.62,7.03,1.3,.24,9,WHITE if s.background.fill.fore_color.rgb==NAVY else GREY)

def bullets(s, items, x=.8,y=1.9,w=11.7,h=4.7,size=24,color=INK):
    tb=s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf=tb.text_frame; tf.clear(); tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.font.name='Aptos'; p.font.size=Pt(size); p.font.color.rgb=color
        p.level=0; p.space_after=Pt(14); p.text='•  '+p.text
    return tb

def notes(s, value):
    tf=s.notes_slide.notes_text_frame; tf.text=value

# 1
s=base('Trust by Design', 'NARSSA · UNISA · SASA | 27 August 2026', True)
text(s,'A Heratio show-and-tell for responsible AI-assisted records management',.72,1.75,11.7,1.2,26,WHITE)
text(s,'Johannes Jurie Pieterse',.72,3.25,7,.45,18,GOLD,True)
text(s,'The Archive & Heritage Group · University of South Africa',.72,3.72,8,.35,14,WHITE)
box(s,.72,5.25,4.7,.78,TEAL,TEAL); text(s,'20 minutes · evidence + live workflow',.94,5.39,4.3,.45,16,WHITE,True)
picture(s,DSAC,9.28,4.42,3.05,2.05,GOLD)
footer(s,1,'0:00–0:45')
notes(s,"""[0:00–0:45] Thank you, Chair, and thank you to NARSSA, UNISA and SASA. I want to use my twenty minutes differently. I will not read a paper to you. I will make one argument, show you the evidence behind it, and then show you the control working inside Heratio.

The argument is simple: records management is an accountability discipline. When AI touches a record, trust cannot be assumed and it cannot be added during an annual audit. Trust has to be designed into the workflow.

By the end, I want you to be able to ask five practical questions of any AI records system: Who remains authorised to decide? What confidence gate is used? What provenance is retained? How is bias tested? And can an auditor reconstruct what happened?""")

# 2
s=base('The question is not “Can AI do it?”','The accountability problem')
box(s,.75,1.85,3.65,3.8,WHITE); text(s,'AI can',1.05,2.12,3,.4,20,TEAL,True)
bullets(s,['classify records','extract metadata','recognise text','find related content'],1.0,2.7,3.0,2.6,20)
text(s,'→',4.63,3.15,.7,.6,34,GOLD,True,PP_ALIGN.CENTER)
box(s,5.35,1.85,7.15,3.8,NAVY,NAVY); text(s,'But can we prove…',5.72,2.12,5.8,.4,20,GOLD,True)
bullets(s,['who made the consequential decision?','why this result was accepted?','what changed—and what did not?','whether error falls unevenly across the archive?'],5.68,2.7,6.35,2.6,20,WHITE)
footer(s,2,'0:45–2:00')
notes(s,"""[0:45–2:00] AI already performs useful work. It can suggest a file-plan class, extract names, recognise typed or handwritten text, and improve retrieval across a large repository. The technical demonstration is no longer the interesting part.

The records-management question is whether the result remains defensible as evidence. If a classifier assigns the wrong retention class, that is not merely a bad search result. It can affect access, retention and disposal. If recognition performs well on modern English correspondence but badly on degraded, handwritten or multilingual material, the system may hide the very voices the archive exists to preserve.

So my test is not whether the model produces an impressive answer. My test is whether we can reconstruct the action: the input, model, confidence, proposed change, human response and final state. That is the difference between AI as a clever feature and AI as a governed records process.""")

# 3
s=base('Five controls around every AI-assisted action','Trust-by-design framework')
labels=[('1','Human authority'),('2','Confidence gates'),('3','Action provenance'),('4','Bias assurance'),('5','Auditability')]
for i,(n,lbl) in enumerate(labels):
    x=.72+i*2.48; box(s,x,2.05,2.18,2.15,WHITE,TEAL); text(s,n,x+.72,2.28,.7,.65,30,TEAL,True,PP_ALIGN.CENTER); text(s,lbl,x+.16,3.22,1.86,.58,17,NAVY,True,PP_ALIGN.CENTER)
text(s,'AI recommends. The accountable professional decides.',1.15,5.23,11.0,.75,25,NAVY,True,PP_ALIGN.CENTER)
picture(s,CONTROL_DIAGRAM,4.25,4.28,4.85,.78,TEAL)
footer(s,3,'2:00–3:20')
notes(s,"""[2:00–3:20] Here is the framework in five controls.

First, human authority: the records professional—not the vendor and not the model—sets the rule and owns consequential decisions. Second, confidence gates: uncertainty must change the route through the workflow. Third, action provenance: retain the model and version, input and output fingerprints, confidence, proposed action and human decision. Fourth, bias assurance: measure performance across meaningful strata—language, period, script, condition and record type—not only as one attractive average. Fifth, auditability: a reviewer must be able to reconstruct the chain after the event.

This is anchored locally: the NARSSA mandate, POPIA, PAIA and the qualities in ISO 15489—authenticity, reliability, integrity and usability. International instruments are useful reference points, but the starting point for South African public institutions must be South African law and recordkeeping obligations.""")

# 4
s=base('Why bias is a records risk—not an abstract AI debate','Evidence from a real archive')
stats=[('454k','records'),('1.35m','digital objects'),('97%','faces returned “unknown”')]
for i,(v,l) in enumerate(stats):
    x=.8+i*4.12; box(s,x,1.95,3.72,2.0,WHITE); text(s,v,x+.18,2.2,3.35,.7,36,TEAL,True,PP_ALIGN.CENTER); text(s,l,x+.18,3.05,3.35,.42,17,GREY,False,PP_ALIGN.CENTER)
picture(s,DOC_A,.8,4.35,2.15,1.72,TEAL); picture(s,DOC_B,3.08,4.35,2.15,1.72,TEAL); picture(s,DOC_C,5.36,4.35,2.15,1.72,TEAL)
box(s,7.72,4.35,4.8,1.72,MINT,MINT); text(s,'AI recognised the already prominent—and under-served almost everyone else.',8.02,4.58,4.2,1.12,20,NAVY,True,PP_ALIGN.CENTER)
footer(s,4,'3:20–5:10')
notes(s,"""[3:20–5:10] Let me make bias concrete. In a large liberation-movement archive, the catalogue contains roughly 454,000 records and 1.35 million digital objects. The descriptive metadata is entirely English, although the underlying record is historically and linguistically much richer.

In one face-recognition run, the system produced 4,906 detections across about 4,800 images. Ninety-seven per cent came back unknown. Two-thirds were below 0.90 confidence, and more than a quarter were below 0.70. The small number recognised were disproportionately already-prominent public figures.

That one run shows why governance matters. Automatic identity assignment would have been indefensible. A confidence score needs a review rule. Identifying a person also introduces a POPIA question. And the uneven result is itself evidence: the machine serves the famous better than the less documented.

The lesson is not that we must reject AI. It is that uncertainty and unevenness must become visible workflow states—not silently written metadata.""")

# 5
s=base('What I will show you in Heratio','Show-and-tell map')
steps=[('1','Open a record'),('2','Ask AI for a suggestion'),('3','Inspect confidence + provenance'),('4','Review, correct, approve'),('5','Reconstruct the audit trail')]
for i,(n,lbl) in enumerate(steps):
    y=1.86+i*.92; text(s,n,.86,y,.48,.48,18,WHITE,True,PP_ALIGN.CENTER); circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(.82),Inches(y-.02),Inches(.55),Inches(.55)); circ.fill.solid(); circ.fill.fore_color.rgb=TEAL; circ.line.fill.background();
    text(s,lbl,1.62,y-.02,7.4,.55,22,NAVY,i in (2,3))
text(s,'The point is the control—not the magic trick.',8.9,2.3,3.2,1.5,24,TEAL,True,PP_ALIGN.CENTER)
box(s,8.72,4.15,3.65,.85,NAVY,NAVY); text(s,'LIVE DEMO · 6 MIN',8.93,4.33,3.25,.42,16,WHITE,True,PP_ALIGN.CENTER)
picture(s,DSAC,9.14,5.22,2.78,1.2,TEAL)
footer(s,5,'5:10–5:45')
notes(s,"""[5:10–5:45] I am now going into Heratio. I will use one safe demonstration record. Watch the sequence rather than the content: open the record; ask for an AI suggestion; inspect the confidence and provenance; make a human decision; then reconstruct the event in the audit trail.

[ACTION] Leave the slideshow and open the prepared Heratio browser tab. Keep this slide available as the fallback map. Do not improvise a different record. Avoid displaying personal or restricted information.""")

# 6 demo
s=base('Show: the record before AI touches it','Heratio live · step 1',True)
box(s,.8,1.8,11.75,3.65,WHITE,WHITE); text(s,'RECORD',1.05,2.08,1.4,.35,11,TEAL,True); text(s,'Original content + existing metadata',1.05,2.55,6.6,.55,25,NAVY,True)
text(s,'No machine assertion should overwrite the source silently.',1.05,3.38,9.9,.62,22,INK)
box(s,9.43,2.35,2.3,1.45,MINT,MINT); text(s,'BASELINE\nSTATE',9.74,2.62,1.7,.8,20,TEAL,True,PP_ALIGN.CENTER)
picture(s,DOC_A,8.26,1.98,3.92,3.18,TEAL)
text(s,'Fallback slide: narrate this if connectivity fails.',.82,6.25,8,.35,13,GOLD,True)
footer(s,6,'5:45–6:45')
notes(s,"""[5:45–6:45 — LIVE DEMO] First I open the record and pause before running anything.

[ACTION] Point out the stable identifier, title, level of description, creator/date fields and digital object. If available, show the version/history indicator.

Say: ‘This is the baseline state. The source record and its existing metadata are not an AI prompt to be overwritten. They are evidence. Any machine contribution must remain distinguishable from what existed before.’

Point out the access classification. Say: ‘Governance begins before inference. The system should know whether the content is permitted to leave the institution or whether only an approved local service may process it.’

[FALLBACK] If the site is unavailable, remain on this slide and explain that the white panel represents the immutable baseline. Continue to the next fallback slide without apologising at length.""")

# 7
s=base('Show: AI proposes; it does not decide','Heratio live · step 2',True)
box(s,.75,1.85,5.55,3.65,WHITE); text(s,'AI SUGGESTION',1.02,2.12,2,.35,11,TEAL,True); text(s,'Suggested classification',1.02,2.62,4.7,.45,23,NAVY,True); text(s,'Confidence 0.78',1.02,3.32,2.7,.45,22,GOLD,True); text(s,'Status: awaiting review',1.02,4.08,3.8,.38,17,RED,True)
text(s,'≠',6.44,3.07,.55,.65,35,GOLD,True,PP_ALIGN.CENTER)
box(s,7.15,1.85,5.35,3.65,NAVY,NAVY); text(s,'RECORD STATE',7.48,2.12,2,.35,11,GOLD,True); text(s,'Unchanged',7.48,2.7,4.3,.6,30,WHITE,True); text(s,'until an authorised human acts',7.48,3.55,4.3,.7,20,WHITE)
picture(s,CONFIDENCE_CHART,7.72,4.37,4.18,1.46,GOLD)
footer(s,7,'6:45–8:00')
notes(s,"""[6:45–8:00 — LIVE DEMO] Now I request an AI-assisted suggestion.

[ACTION] Trigger the prepared AI function—classification, entity extraction or description suggestion. Do not choose an operation that can write directly to the record. When the result appears, point to the label ‘suggestion’, the confidence score and the review status.

Say: ‘The crucial design decision is visible here. A generated answer is not yet record metadata. It is a proposed assertion awaiting an authorised decision.’

If the confidence is different from the rehearsal value, use the actual number. Explain the threshold: low confidence routes to mandatory review; higher confidence may reduce priority but should not transfer statutory appraisal or disposal authority to the model.

Point back to the record and say: ‘Notice that the source is unchanged. This separation protects integrity and makes correction ordinary rather than exceptional.’""")

# 8
s=base('Show: provenance travels with the suggestion','Heratio live · step 3',True)
fields=[('MODEL','service · model · version'),('INPUT','record/object fingerprint'),('OUTPUT','suggestion fingerprint'),('CONTEXT','confidence · rule · timestamp'),('REVIEW','reviewer · decision · reason')]
for i,(a,b) in enumerate(fields):
    y=1.75+i*.92; box(s,.82,y,11.7,.68,RGBColor(22,52,67),RGBColor(46,79,91)); text(s,a,1.05,y+.12,1.55,.36,12,GOLD,True); text(s,b,2.65,y+.1,8.9,.4,18,WHITE)
footer(s,8,'8:00–9:15')
notes(s,"""[8:00–9:15 — LIVE DEMO] Open the provenance or AI event detail.

[ACTION] Point, in order, to the service/model identity, version if exposed, timestamp, input and output references or hashes, confidence, and review state.

Say: ‘A screenshot of an AI answer is not provenance. For audit purposes we need enough information to identify what service acted, on what input, what it proposed, under which rule, and what happened next.’

Do not claim that a hash proves the model was correct. Say: ‘Tamper evidence protects the history of the action; it does not certify the truth of the output. Accuracy still requires testing and human judgement.’

Also note data sovereignty if the demo uses the locally hosted service: ‘The processing location and approved provider are governance attributes, not hidden infrastructure details.’""")

# 9
s=base('Show: the human decision becomes part of the record','Heratio live · step 4',True)
for i,(lab,col) in enumerate([('REJECT',RED),('CORRECT',GOLD),('APPROVE',TEAL)]):
    x=.9+i*4.1; box(s,x,2.0,3.55,1.15,col,col); text(s,lab,x+.2,2.28,3.15,.48,22,WHITE,True,PP_ALIGN.CENTER)
text(s,'Decision + reason + reviewer + time',1.1,4.05,11.1,.65,28,WHITE,True,PP_ALIGN.CENTER)
text(s,'Human judgement is not friction. It is the accountable act.',1.0,5.16,11.3,.65,23,GOLD,True,PP_ALIGN.CENTER)
footer(s,9,'9:15–10:45')
notes(s,"""[9:15–10:45 — LIVE DEMO] Now I make the accountable decision.

[ACTION] Deliberately correct one part of the suggestion rather than simply approving it. Enter a short reason, for example: ‘Context in the record indicates correspondence, not policy.’ Save the review.

Say: ‘Correction is the most revealing path because it proves the interface is not ceremonial. The reviewer can reject, amend or approve. The reason is attributed and retained.’

Point out the identity and timestamp. If dual approval exists for sensitive actions, mention it without opening a new workflow. State the hard boundary clearly: ‘AI may recommend appraisal or retention information, but statutory appraisal and disposal remain human-authorised actions.’

Then open the record’s updated metadata. Say: ‘Only after the decision does the accepted value become part of the operational record—and the machine contribution remains traceable.’""")

#10
s=base('Tell: reconstruct the chain after the event','Heratio live · step 5')
chain=[('1','Baseline'),('2','Inference'),('3','Suggestion'),('4','Review'),('5','Record update')]
for i,(n,lbl) in enumerate(chain):
    x=.75+i*2.48; c=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+.58),Inches(2.1),Inches(.75),Inches(.75)); c.fill.solid(); c.fill.fore_color.rgb=TEAL; c.line.fill.background(); text(s,n,x+.72,2.25,.46,.34,16,WHITE,True,PP_ALIGN.CENTER); text(s,lbl,x,3.12,1.9,.45,17,NAVY,True,PP_ALIGN.CENTER)
    if i<4: text(s,'→',x+1.78,2.18,.62,.45,24,GOLD,True,PP_ALIGN.CENTER)
box(s,1.25,4.45,10.8,1.05,MINT,MINT); text(s,'Can an auditor explain what happened without trusting our memory?',1.55,4.7,10.2,.5,23,NAVY,True,PP_ALIGN.CENTER)
picture(s,CONTROL_DIAGRAM,4.2,5.72,4.95,.82,TEAL)
footer(s,10,'10:45–12:00')
notes(s,"""[10:45–12:00 — LIVE DEMO] Finally, open the audit or activity trail and reconstruct the chain.

[ACTION] Read the sequence, not every field: baseline record; inference request; generated suggestion; human correction and reason; final update.

Say: ‘This is the practical test of trust by design. Months later, an auditor should not have to trust my recollection or the vendor’s marketing. The system should tell the story.’

Return to the slideshow after showing the chain.

[TRANSITION] ‘What you have just seen is deliberately modest. Heratio did not replace the records manager. It made a suggestion, exposed uncertainty, retained provenance, required judgement and preserved the decision. That modesty is precisely what makes the workflow defensible.’""")

#11
s=base('What this changes in institutional practice','From demo to operating model')
items=[('Policy','Define prohibited and human-only decisions.'),('Workflow','Route by confidence and sensitivity.'),('Assurance','Test by language, era, script and condition.'),('Evidence','Retain inference, review and correction history.')]
for i,(a,b) in enumerate(items):
    y=1.75+i*1.03; text(s,a,.82,y,2.0,.48,18,TEAL,True); text(s,b,2.75,y,9.3,.48,22,NAVY)
box(s,.82,6.05,11.65,.55,NAVY,NAVY); text(s,'Start with one bounded workflow. Measure. Review. Expand deliberately.',1.1,6.12,11.0,.38,18,WHITE,True,PP_ALIGN.CENTER)
footer(s,11,'12:00–14:20')
notes(s,"""[12:00–14:20] The demonstration only matters if it changes operating practice.

At policy level, define decisions that AI may never make alone: appraisal, disposal, access restriction and identity assignment in sensitive contexts. At workflow level, route work using both confidence and consequence. A high-confidence low-risk suggestion is not the same as a high-confidence disposal recommendation.

At assurance level, stop reporting only one overall accuracy score. Test by meaningful strata: language, period, handwriting or print, physical condition, office of origin and record type. A model can have an acceptable average while systematically failing a historically important subset.

At evidence level, retain corrections. A corrected suggestion is not an embarrassment to delete; it is assurance evidence and future training signal.

My practical recommendation is to begin with one bounded workflow—such as metadata suggestion—establish the baseline, set the review rule, measure the error distribution, and expand only when the control performs as intended.""")

#12
s=base('A minimum assurance checklist','What to ask before procurement or deployment')
bullets(s,['Is every machine assertion visibly distinguishable from human-authored metadata?','Can the records manager set review gates and override the result?','Are model, input, output, confidence and review captured together?','Is performance measured across local languages and record conditions?','Can the institution export the audit history and retain control of its data?'],.85,1.72,11.65,4.85,20)
footer(s,12,'14:20–16:50')
notes(s,"""[14:20–16:50] Here is the checklist I would take into a procurement meeting tomorrow.

First: can users distinguish machine assertions from human-authored metadata? Second: can the accountable professional configure gates and override results? Third: are the model identity, input, output, confidence and review joined as one provenance chain? Fourth: has performance been measured on our languages, periods and record conditions—not only on a vendor benchmark? Fifth: can the institution export the audit history and retain sovereignty over its records and sensitive data?

If a system cannot answer these questions, the problem is not merely incomplete functionality. It is an assurance gap.

Also ask what happens when the model changes. Version changes can change outputs. A defensible programme therefore treats model updates like controlled changes: document, test, approve, monitor and retain the earlier context for actions already taken.""")

#13
s=base('Three things to remember','Close',True)
for i,(n,t) in enumerate([('1','AI output is a proposal—not a recordkeeping decision.'),('2','Uncertainty must alter the workflow.'),('3','Trust is demonstrated through provenance and accountable review.')]):
    y=1.65+i*1.35; text(s,n,.85,y,.62,.62,26,GOLD,True,PP_ALIGN.CENTER); text(s,t,1.75,y,10.6,.7,24,WHITE,True)
box(s,.85,5.85,4.5,.68,TEAL,TEAL); text(s,'Heratio: trust made visible',1.08,5.96,4.05,.42,17,WHITE,True,PP_ALIGN.CENTER)
picture(s,DSAC,10.62,5.08,1.62,1.28,GOLD)
footer(s,13,'16:50–18:40')
notes(s,"""[16:50–18:40] I will close with three points.

First, AI output is a proposal, not a recordkeeping decision. The professional mandate remains human. Second, uncertainty must change the workflow. A confidence score displayed in small print is not a control unless it routes, pauses or escalates an action. Third, trust is not a claim made by the system supplier. It is demonstrated through provenance, review and an audit trail that another person can inspect.

Heratio is the show-and-tell, but the framework is deliberately platform-independent. Any institution can apply these controls to its own architecture.

South African archives hold contested, multilingual and historically uneven records. We should use AI to make them more usable—but never by making its interventions invisible. Responsible adoption is not slower innovation. It is innovation that can survive scrutiny.""")

#14
s=base('Thank you','Questions',True)
text(s,'Johannes Jurie Pieterse',.75,2.0,6,.5,26,WHITE,True)
text(s,'johan@theahg.co.za',.75,2.62,6,.42,18,GOLD)
text(s,'ORCID 0000-0002-1670-2416',.75,3.13,6,.42,16,WHITE)
text(s,'Prompt for discussion:',7.2,2.02,4.7,.4,16,GOLD,True)
text(s,'Which AI-assisted decision in your institution most urgently needs a visible human review gate?',7.2,2.62,4.9,1.55,23,WHITE,True)
picture(s,DSAC,.78,4.3,3.25,2.15,GOLD)
footer(s,14,'18:40–20:00')
notes(s,"""[18:40–20:00] Thank you.

If the Chair invites immediate discussion, ask: ‘Which AI-assisted decision in your institution most urgently needs a visible human review gate?’

Likely Q&A:
• Is the 97% unknown rate proof of racial bias? Answer: No. It is evidence of highly uneven usefulness and representation risk in that collection; causal bias claims require controlled evaluation.
• Does provenance prove correctness? No. It proves traceability and supports accountability; accuracy requires validation.
• Can AI ever auto-approve? Only in bounded, low-consequence workflows under an authorised policy, monitored thresholds and reversible actions. Appraisal and disposal remain human.
• Why local hosting? It can support sovereignty and confidentiality, but location alone is not governance. Access, logs, models, retention and security still require controls.
• Is Heratio the framework? No. Heratio demonstrates it; the five controls are platform-independent.

[TIME] Finish speaking by 18:40 to preserve at least 80 seconds for the chair or one question.""")

prs.core_properties.title='Trust by Design — Heratio Show-and-Tell (20-minute draft)'
prs.core_properties.subject='NARSSA/UNISA/SASA 2nd National Records Management Conference 2026'
prs.core_properties.author='Johannes Jurie Pieterse'
prs.core_properties.comments='Draft generated 22 August 2026; includes full speaker notes and live-demo fallback slides.'
prs.save(OUT)
print(OUT)
